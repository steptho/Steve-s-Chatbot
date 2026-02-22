import streamlit as st
import os
import json
from datetime import datetime
from openai import OpenAI
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation  # Added for PowerPoint support
from PIL import Image
import base64

# ----------------------------
# HELPER FUNCTIONS
# ----------------------------

def generate_chat_title(text):
    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Generate a short professional conversation title (max 6 words). Respond ONLY with the title text."},
                {"role": "user", "content": text[:2000]}
            ],
        )
        title = response.choices[0].message.content.strip().replace('"', '').replace(':', '').replace('/', '')
        return title
    except:
        return datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

def rename_chat_file(old_id, new_title):
    old_path = os.path.join(CHAT_DIR, f"{old_id}.json")
    timestamp = datetime.now().strftime("%H%M")
    new_id = f"{new_title}_{timestamp}"
    new_path = os.path.join(CHAT_DIR, f"{new_id}.json")
    
    if os.path.exists(old_path):
        os.rename(old_path, new_path)
        st.session_state.current_chat_id = new_id
        return new_id
    return old_id

def analyse_image(uploaded_file):
    bytes_data = uploaded_file.read()
    base64_image = base64.b64encode(bytes_data).decode("utf-8")
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{
            "role": "user",
            "content": [
                {"type": "text", "text": "Provide a detailed professional analysis of this image."},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}},
            ],
        }],
    )
    return response.choices[0].message.content

# -------------------------------------------------
# CONFIG & SESSION INIT
# -------------------------------------------------
st.set_page_config(page_title="Steve's Chatbot", layout="wide")
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

CHAT_DIR = "saved_chats"
os.makedirs(CHAT_DIR, exist_ok=True)

for key in ["messages", "current_chat_id", "uploaded_files"]:
    if key not in st.session_state:
        st.session_state[key] = [] if key != "current_chat_id" else None

# -------------------------------------------------
# SAVE / LOAD
# -------------------------------------------------
def save_chat(chat_id):
    if chat_id:
        with open(os.path.join(CHAT_DIR, f"{chat_id}.json"), "w") as f:
            json.dump(st.session_state.messages, f, indent=2)

def load_chat(file_name):
    with open(os.path.join(CHAT_DIR, file_name), "r") as f:
        st.session_state.messages = json.load(f)
    st.session_state.current_chat_id = file_name.replace(".json", "")

# -------------------------------------------------
# SIDEBAR
# -------------------------------------------------
st.sidebar.title("🗂 Conversations")

if st.sidebar.button("➕ New Conversation"):
    st.session_state.messages = []
    st.session_state.current_chat_id = None
    st.rerun()

st.sidebar.divider()

# List saved chats
chat_files = sorted(
    [f for f in os.listdir(CHAT_DIR) if f.endswith(".json")],
    key=lambda x: os.path.getmtime(os.path.join(CHAT_DIR, x)),
    reverse=True
)

for file in chat_files:
    display_name = file.replace(".json", "").replace("_", " ")
    if st.sidebar.button(display_name, key=file):
        load_chat(file)
        st.rerun()

st.sidebar.divider()
st.sidebar.header("📁 Upload Files")

uploaded_files = st.sidebar.file_uploader(
    "Upload PDF, Word, PPT, Excel, Images, CSV",
    type=["pdf", "docx", "pptx", "xlsx", "csv", "png", "jpg", "jpeg"],
    accept_multiple_files=True
)

if uploaded_files:
    st.session_state.uploaded_files = uploaded_files
    if st.sidebar.button("🔍 Analyse Uploaded Files"):
        analysis_text = ""
        with st.spinner("Analyzing files..."):
            for file in uploaded_files:
                # Images
                if file.type.startswith("image"):
                    analysis_text += f"\n\n[IMAGE: {file.name}]\n" + analyse_image(file)
                # PDFs
                elif file.name.endswith(".pdf"):
                    reader = PdfReader(file)
                    for page in reader.pages:
                        analysis_text += page.extract_text() + "\n"
                # Word
                elif file.name.endswith(".docx"):
                    doc = Document(file)
                    for para in doc.paragraphs:
                        analysis_text += para.text + "\n"
                # PowerPoint
                elif file.name.endswith(".pptx"):
                    prs = Presentation(file)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                analysis_text += shape.text + "\n"
                # Data (Excel/CSV)
                elif file.name.endswith((".xlsx", ".csv")):
                    df = pd.read_excel(file) if file.name.endswith(".xlsx") else pd.read_csv(file)
                    analysis_text += f"\nData Summary for {file.name}:\n" + df.head().to_string()

        # Send analysis to LLM
        if analysis_text:
            if st.session_state.current_chat_id is None:
                st.session_state.current_chat_id = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            
            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "system", "content": "Analyze these files."}, {"role": "user", "content": analysis_text[:15000]}]
            )
            reply = response.choices[0].message.content
            st.session_state.messages.append({"role": "assistant", "content": reply})
            
            # Generate title if it's a new chat
            new_title = generate_chat_title(analysis_text[:1000])
            st.session_state.current_chat_id = rename_chat_file(st.session_state.current_chat_id, new_title)
            save_chat(st.session_state.current_chat_id)
            st.rerun()

# -------------------------------------------------
# MAIN CHAT
# -------------------------------------------------
st.title("💬 Steve's Professional Chatbot")

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if prompt := st.chat_input("Ask something..."):
    is_new = st.session_state.current_chat_id is None
    if is_new:
        st.session_state.current_chat_id = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

    st.session_state.messages.append({"role": "user", "content": prompt})
    with st.chat_message("user"): st.markdown(prompt)

    response = client.chat.completions.create(model="gpt-4o-mini", messages=st.session_state.messages)
    reply = response.choices[0].message.content
    st.session_state.messages.append({"role": "assistant", "content": reply})
    with st.chat_message("assistant"): st.markdown(reply)

    if is_new:
        new_title = generate_chat_title(prompt)
        st.session_state.current_chat_id = rename_chat_file(st.session_state.current_chat_id, new_title)
    
    save_chat(st.session_state.current_chat_id)
    st.rerun()